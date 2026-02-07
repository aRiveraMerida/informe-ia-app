# -*- coding: utf-8 -*-
"""Tests para PPTXReportGenerator."""
import pytest
from pptx import Presentation
from io import BytesIO

from modules.pptx_generator import PPTXReportGenerator


class TestPPTXGeneration:
    def test_generates_valid_pptx_bytes(self):
        gen = PPTXReportGenerator(
            client_name="Test Client",
            period="01/01/2025 - 31/01/2025",
        )
        result = gen.generate(
            analysis_text="## Resumen\nEsto es un test.\n## Detalle\nMás texto.",
            metadata={"total_records": 100},
        )

        assert isinstance(result, bytes)
        assert len(result) > 0

        # Verificar que es un PPTX válido
        prs = Presentation(BytesIO(result))
        assert len(prs.slides) >= 1

    def test_title_slide_present(self):
        gen = PPTXReportGenerator(
            client_name="Movimer",
            period="Febrero 2025",
            report_title="Informe Ejecutivo",
        )
        result = gen.generate(
            analysis_text="## Sección 1\nTexto de prueba.",
            metadata={},
        )

        prs = Presentation(BytesIO(result))
        # Al menos portada + 1 sección
        assert len(prs.slides) >= 2

    def test_closing_slide_with_cost(self):
        gen = PPTXReportGenerator(client_name="Test", period="2025")
        result = gen.generate(
            analysis_text="## Test\nContenido.",
            metadata={},
            cost_info={
                "total_cost_usd": 0.0123,
                "total_input_tokens": 1000,
                "total_output_tokens": 2000,
            },
        )

        prs = Presentation(BytesIO(result))
        # Portada + sección + closing
        assert len(prs.slides) >= 3

    def test_markdown_simplification(self):
        text = "**negrita** y *cursiva* y ## heading"
        result = PPTXReportGenerator._simplify_markdown(text)
        assert "**" not in result
        assert "*" not in result

    def test_empty_analysis_still_works(self):
        gen = PPTXReportGenerator(client_name="X", period="Y")
        result = gen.generate(analysis_text="", metadata={})

        assert isinstance(result, bytes)
        prs = Presentation(BytesIO(result))
        assert len(prs.slides) >= 1  # Al menos la portada
