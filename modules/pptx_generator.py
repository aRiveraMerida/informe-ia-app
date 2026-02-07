# -*- coding: utf-8 -*-
"""
Generador de presentaciones PowerPoint (PPTX) — Estilo Movimer.
Portada, resumen ejecutivo, secciones con gráficos, conclusiones.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from io import BytesIO
from datetime import datetime
from typing import Optional, Dict, Any, List
import re
import logging

logger = logging.getLogger(__name__)

# Colores Movimer
_GREEN = RGBColor(112, 174, 0)
_GREEN_DARK = RGBColor(61, 107, 0)
_GRAY = RGBColor(87, 87, 87)
_BLACK = RGBColor(0, 0, 0)
_WHITE = RGBColor(255, 255, 255)


class PPTXReportGenerator:
    """Genera presentaciones PowerPoint profesionales estilo Movimer."""

    def __init__(
        self,
        client_name: str,
        period: str,
        report_title: str = "Informe Ejecutivo",
    ):
        self.client_name = client_name
        self.period = period
        self.report_title = report_title
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)

    def generate(
        self,
        analysis_text: str,
        metadata: Dict[str, Any],
        cost_info: Optional[Dict[str, Any]] = None,
        report_sections: Optional[List] = None,
    ) -> bytes:
        """Genera PPTX completo y devuelve bytes."""
        self._add_title_slide()

        if report_sections:
            for section in report_sections:
                self._add_section_slide(section)
        else:
            self._add_markdown_slides(analysis_text)

        if cost_info:
            self._add_closing_slide(cost_info)

        buffer = BytesIO()
        self.prs.save(buffer)
        buffer.seek(0)
        return buffer.getvalue()

    # ─── Slides ───────────────────────────────────────────────────────────────

    def _add_title_slide(self):
        """Slide de portada."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # blank
        self._set_slide_bg(slide, _WHITE)

        # Barra verde superior
        self._add_shape_rect(slide, 0, 0, self.prs.slide_width, Inches(0.15), _GREEN)

        # Título
        self._add_textbox(
            slide, Inches(1), Inches(2), Inches(11), Inches(1.5),
            self.report_title.upper(),
            font_size=40, bold=True, color=_BLACK, alignment=PP_ALIGN.CENTER,
        )

        # Línea verde
        self._add_shape_rect(
            slide, Inches(4.5), Inches(3.5), Inches(4), Inches(0.05), _GREEN,
        )

        # Cliente
        self._add_textbox(
            slide, Inches(1), Inches(3.8), Inches(11), Inches(0.8),
            self.client_name,
            font_size=24, bold=True, color=_GREEN_DARK, alignment=PP_ALIGN.CENTER,
        )

        # Periodo y fecha
        fecha = datetime.now().strftime("%d/%m/%Y")
        self._add_textbox(
            slide, Inches(1), Inches(5), Inches(11), Inches(0.6),
            f"Periodo: {self.period}  |  Fecha: {fecha}",
            font_size=14, color=_GRAY, alignment=PP_ALIGN.CENTER,
        )

        # Barra verde inferior
        self._add_shape_rect(
            slide, 0, self.prs.slide_height - Inches(0.15),
            self.prs.slide_width, Inches(0.15), _GREEN,
        )

    def _add_section_slide(self, section):
        """Slide para una sección del informe con su texto y gráficos."""
        if not section.heading and not section.text:
            return

        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_slide_bg(slide, _WHITE)
        self._add_shape_rect(slide, 0, 0, self.prs.slide_width, Inches(0.08), _GREEN)

        y_offset = Inches(0.4)

        # Heading
        if section.heading:
            self._add_textbox(
                slide, Inches(0.8), y_offset, Inches(11.5), Inches(0.6),
                section.heading,
                font_size=24, bold=True, color=_GREEN_DARK,
            )
            y_offset += Inches(0.7)

        # Texto (simplificado: extraer bullets y párrafos)
        if section.text:
            clean_text = self._simplify_markdown(section.text)
            # Limitar texto para que quepa en el slide
            max_chars = 800 if not section.chart_images else 400
            if len(clean_text) > max_chars:
                clean_text = clean_text[:max_chars] + "..."

            text_height = Inches(2.5) if not section.chart_images else Inches(1.8)
            self._add_textbox(
                slide, Inches(0.8), y_offset, Inches(11.5), text_height,
                clean_text,
                font_size=12, color=_BLACK,
            )
            y_offset += text_height + Inches(0.2)

        # Gráficos (máximo 2 por slide)
        if section.chart_images:
            chart_count = min(2, len(section.chart_images))
            chart_width = Inches(5.5) if chart_count == 2 else Inches(8)
            for i, (title, png_bytes) in enumerate(section.chart_images[:chart_count]):
                try:
                    img_stream = BytesIO(png_bytes)
                    x = Inches(0.8) + (Inches(6) * i) if chart_count == 2 else Inches(2.5)
                    slide.shapes.add_picture(
                        img_stream, x, y_offset, chart_width, Inches(3),
                    )
                except Exception as e:
                    logger.warning("No se pudo insertar gráfico '%s' en PPTX: %s", title, e)

    def _add_markdown_slides(self, markdown_text: str):
        """Convierte markdown en slides (un slide por sección ##)."""
        sections = re.split(r'^## ', markdown_text, flags=re.MULTILINE)
        for section_text in sections:
            if not section_text.strip():
                continue
            lines = section_text.strip().split("\n", 1)
            heading = lines[0].strip().lstrip("#").strip()
            body = lines[1].strip() if len(lines) > 1 else ""

            slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
            self._set_slide_bg(slide, _WHITE)
            self._add_shape_rect(slide, 0, 0, self.prs.slide_width, Inches(0.08), _GREEN)

            self._add_textbox(
                slide, Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.6),
                heading, font_size=24, bold=True, color=_GREEN_DARK,
            )

            clean_body = self._simplify_markdown(body)
            if len(clean_body) > 1200:
                clean_body = clean_body[:1200] + "..."

            self._add_textbox(
                slide, Inches(0.8), Inches(1.2), Inches(11.5), Inches(5.5),
                clean_body, font_size=12, color=_BLACK,
            )

    def _add_closing_slide(self, cost_info: Dict[str, Any]):
        """Slide final con metadatos de generación."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_slide_bg(slide, _WHITE)
        self._add_shape_rect(slide, 0, 0, self.prs.slide_width, Inches(0.08), _GREEN)

        fecha = datetime.now().strftime("%d/%m/%Y %H:%M")
        cost = cost_info.get("total_cost_usd", 0)
        tokens_in = cost_info.get("total_input_tokens", 0)
        tokens_out = cost_info.get("total_output_tokens", 0)

        self._add_textbox(
            slide, Inches(1), Inches(2.5), Inches(11), Inches(1),
            "Informe generado con IA",
            font_size=28, bold=True, color=_GREEN_DARK, alignment=PP_ALIGN.CENTER,
        )

        meta = (
            f"Fecha: {fecha}  |  Tokens: {tokens_in:,} + {tokens_out:,}  |  "
            f"Coste: ${cost:.4f} USD"
        )
        self._add_textbox(
            slide, Inches(1), Inches(4), Inches(11), Inches(0.5),
            meta, font_size=12, color=_GRAY, alignment=PP_ALIGN.CENTER,
        )

        self._add_shape_rect(
            slide, 0, self.prs.slide_height - Inches(0.15),
            self.prs.slide_width, Inches(0.15), _GREEN,
        )

    # ─── Utilidades ───────────────────────────────────────────────────────────

    @staticmethod
    def _simplify_markdown(text: str) -> str:
        """Convierte markdown a texto plano simplificado para slides."""
        text = re.sub(r'^#{1,4}\s+', '', text, flags=re.MULTILINE)
        text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)
        text = re.sub(r'__(.+?)__', r'\1', text)
        text = re.sub(r'\*(.+?)\*', r'\1', text)
        text = re.sub(r'_(.+?)_', r'\1', text)
        text = re.sub(r'^---+$', '', text, flags=re.MULTILINE)
        # Convertir tablas markdown a texto tabulado simple
        text = re.sub(r'\|[-:]+\|[-:| ]+\|', '', text)
        text = text.replace("|", "  ")
        # Limpiar líneas vacías múltiples
        text = re.sub(r'\n{3,}', '\n\n', text)
        return text.strip()

    @staticmethod
    def _add_textbox(slide, left, top, width, height, text,
                     font_size=12, bold=False, color=_BLACK, alignment=PP_ALIGN.LEFT):
        """Añade un textbox con formato."""
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.alignment = alignment

    @staticmethod
    def _add_shape_rect(slide, left, top, width, height, fill_color):
        """Añade un rectángulo de color sólido."""
        from pptx.util import Emu as _Emu
        shape = slide.shapes.add_shape(
            1, left, top, width, height,  # 1 = MSO_SHAPE.RECTANGLE
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.fill.background()

    @staticmethod
    def _set_slide_bg(slide, color):
        """Establece el color de fondo de un slide."""
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color
